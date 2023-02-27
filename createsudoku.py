import io
import random
import string
import aspose.slides as slides2
import aspose.pydrawing as drawing
from pptx.util import Pt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import os
import PySimpleGUI as gui
import constants
import requests
import sys
from pathlib import Path
from subprocess import check_output, CalledProcessError, STDOUT
import shlex

THEME = constants.THEME


class Classcreateasudoku:

    def system_call(self,command):
        """ 
        params:
            command: string, ex. `"ls -l"`
        returns: output, success
        """
        command = shlex.split(command)
        try:
            output = check_output(command,shell=True, stderr=STDOUT).decode()
            success = True 
        except CalledProcessError as e:
            output = e.output.decode()
            success = False
        return output, success

    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _set_top_cell_border(self,cell, border_color="000000", border_width='90800'):
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        lnT = self.SubElement(
            tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnT, 'a:solidFill')
        srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
        return cell

    def _set_bottom_cell_border(self,cell, border_color="000000", border_width='90800'):
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        lnB = self.SubElement(
            tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnB, 'a:solidFill')
        srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
        return cell

    def _set_left_cell_border(self,cell, border_color="000000", border_width='90800'):
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        lnL = self.SubElement(
            tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnL, 'a:solidFill')
        srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
        return cell

    def _set_right_cell_border(self,cell, border_color="000000", border_width='90800'):
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        lnR = self.SubElement(
            tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = self.SubElement(lnR, 'a:solidFill')
        srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
        return cell

    def move_slide(self,old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

    
    def collapse(self, layout, key, visible):
        """
        Helper function that creates a Column that can be later made hidden, thus appearing "collapsed"
        :param layout: The layout for the section
        :param key: Key used to make this section visible / invisible
        :param visible: visible determines if section is rendered visible or invisible on initialization
        :return: A pinned column that can be placed directly into your layout
        :rtype: sg.pin
        """
        return gui.pin(gui.Column(layout, key=key, visible=visible, pad=(0, 0)))

    def __init__(self):
        self.closed = False
        self.fonts = gui.Text.fonts_installed_list()
        self.difficultylevel = ["simple","easy","intermediate","expert"]

        self.bpicon_base64 = b"iVBORw0KGgoAAAANSUhEUgAABAAAAAQACAIAAADwf7zUAAAACXBIWXMAAC4jAAAuIwF4pT92AABV70lEQVR4nO3de4CXc/7w/5rppDI66WQ0UtF2UFQoCYmsuq2QSijk9HVai91t0TrknHOWUCyJjtqVhJYoFUXRQSjpYFR0UKppmsPvD/dv7733jj4zzcz7uj6fx+Ov7199nztm3p/363pfn+sqVw4AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAOKlfOgASFSbNm06d+7cqlWrzMzMBg0a1KxZMyMjo2rVqvvtt1/58n6TI2HHjh3lypXLzc3Ny8srKCjYtm3btm3btm7dumXLls2bN2/evHnTpk1ff/314sWLP/3009CxUDRZWVknnXRS69atDz744IMOOqhWrVoZGRnVq1evUKFCenp65cqVQwcmlby8vNzc3J//78LCwp07d+bl5e3atSs3N3fHjh3bt2/fvn37tm3bfl5bvv/++/Xr169evXrBggWbN28OWw6xYNtEpHXu3LlXr17t27dv0aJFnTp1QudQYnJzczds2LB+/frs7Ow1a9YsWrTo/fffX7p0aegu+L+0adPmnHPO6dChQ+vWrRs2bBg6h73btm3b5s2b169fv27durVr165evXrp0qUzZszYunVr6DSIEAMAUXT66aeff/75xx9/fGZmZugWys7mzZu/+eabFStWLF269IMPPnjrrbdCF5GiOnfufNFFF3Xp0qVJkyYOGJNAYWHh999/v2bNmm+++WbZsmULFy586623jASkMusaEZKVlXXttdf26tWrcePGoVsI78cff/z8888XLlz4wQcfTJgwIScnJ3QRSa527drXX399r169WrRoEbqF0lVQULBq1aqvvvpqyZIl8+fP/8c//rF9+/bQUVB2DABEwjHHHHPTTTf99re/rVq1augWoig3N3fZsmWzZ8+eMmXK66+/HjqHZNOqVas//vGPZ5555v777x+6hQByc3NXrFixePHiDz/8cPLkyStWrAhdBKXLAEBg7du3v+WWW3r06FGhQoXQLcTDDz/8MG/evHffffell17Kzs4OnUO8HXbYYUOHDj3jjDN8i5efFRYWrly5cuHChTNnznzllVfWrVsXughKngGAYOrWrfvAAw/07du3UqVKoVuIpby8vI8//viNN94YOXLk2rVrQ+cQM9WqVRs2bNiFF17o4JFfkpeXt3Tp0jlz5vzzn/+cOnVq6BwoMQYAwvj9739/8803e7APJWL37t3z5s2bOnXqM888s2HDhtA5xMDll18+ZMgQD/YhcevWrZszZ87rr7/+/PPP5+fnh86BfWIAoKwddthhTz/99AknnBA6hCSUk5Pz/vvvv/DCCy+99FLoFiIqKytrxIgR3bt3Dx1CXG3ZsmXu3LmvvfbaM888s3v37tA5UBwGAMrUZZddds8999SqVSt0CElu1apVkydPHj58+PLly0O3ECEXXnjhAw88ULdu3dAhJIMtW7a89957Y8aMGTduXOgWKBoDAGUkLS3t2WefHThwoIdqU2Zyc3NnzZr11FNPjR8/PnQL4T3xxBOXX355enp66BCSTXZ29uuvvz5ixIiPP/44dAskxFaMslC/fv3Jkycfc8wxoUNIUYsWLRo5cuTjjz9eUFAQuoUAatasOXny5C5duoQOIZnl5+d/8sknY8eOfeSRR3xJgIgzAFDq2rVrN27cuEMPPTR0CKluzZo1o0ePvvfee70BNKW0aNFi4sSJzZs3Dx1Cqvj+++9fe+21Rx55ZNGiRaFbYM8MAJSurl27jhkzpl69eqFD4H/bsmXLCy+8cOuttxoDUkHHjh0nTJjgaT+Uvby8vBkzZjzxxBOTJ08O3QL/zQBAKerWrdvYsWN95ZcI2rhx43PPPXfrrbfm5OSEbqG0dOvWbfTo0S5AENZnn302cuTIxx57LHQI/B8GAErLySefPG7cOLt/omz9+vWjRo269dZb3bCbfDp27Dh58mQP/CEiVq5cOWrUqHvuucdqQxR4GAKlol27dhMnTjzwwANDh8CvqV69+vHHH3/xxRfn5eV99NFHoXMoMa1atXr99ddd+yc6atas2bVr10suuWT//fefNWtWYWFh6CJSmgGAkle/fv0333wzMzMzdAgkJCMj47e//W2PHj1Wrlz59ddfh85hX9WuXfvNN9885JBDQofAf8vIyDjxxBONAQRnAKCEpaWlvf32261atQodAkVz0EEH9e/f/4gjjpg3b96WLVtC51B8b7755lFHHRW6An7Rz2PAhRdeWFhY+OGHH4bOIRUZAChho0aN6tmzZ+gKKI60tLQWLVpcdNFF1apVe/fdd0PnUBxPPvnk2WefHboC9q5GjRq//e1ve/fuvXnz5sWLF4fOIbX4EjAl6Yorrvjb3/7mXb8kgQULFlx77bWzZs0KHUIRDBw4cOTIkWlpaaFDoAgKCwvffffdwYMH+yYSZcZGjRLTvHnzDz74wGN/SBq7du169tlnr7vuOk/tiIUmTZrMmTPHsweIqdzc3DFjxvzhD3/YvHlz6BaSn1uAKDHetUmSqVChwtFHH92vX79Vq1Z98cUXoXPYi/Hjx/v2EfGVnp7etm3bAQMG5OXl+WIApc05KSXjD3/4Q5cuXUJXQMlr1qzZxIkTn3rqqfR0V0yi6+qrrz755JNDV8C+ql+//iOPPDJ79uz27duHbiGZuQWIElC3bt0lS5bUqVMndAiUoiVLllx22WWzZ88OHcJ/q1mz5rJly7zzi2SyY8eOESNG3HjjjQUFBaFbSEIuaFECnnrqqWOPPTZ0BZSuunXr9u3bt6CgwDeDo+aJJ57o3Llz6AooSRUrVuzYsWPv3r2XLl36zTffhM4h2TgBYF+1b9/+gw8+qFSpUugQKCOvv/76BRdc4It6EdGqVat58+ZVqVIldAiUil27do0YMeL66693FEAJ8h0A9tVf//pXu39SSo8ePebNm3fCCSeEDqFcuXLlhg4davdPEqtcufK11167cOHCY445JnQLycMtQOyTjh073n333Z66TaqpVatW7969d+zYMXfu3NAtKa19+/b33nuv72eT9OrVq9evX7/09PT3338/dAvJwL6NfXLDDTf46CU1VatW7aGHHvLaqbBuvPHGihUrhq6AslCtWrU77rjj7bffbtiwYegWYs93ACi+xo0bL1261OE7KW7WrFm9e/det25d6JCU07Bhw6+++qpq1aqhQ6BMrV279qqrrvrnP/8ZOoQYc+GK4rv22mvt/qFz586zZs1q165d6JCU8/vf/97unxSUmZk5fvz4u+66K3QIMeYEgOL7+uuvGzduHLoCImHjxo2DBg2aPHly6JAUsnz58iZNmoSugGBef/318847b+vWraFDiB8nABTT7373O7t/+LfatWuPGTPmiiuuCB2SKk4//XS7f1Jcjx495s+f753BFIMBgGLq06dP6ASIlv3222/48OFDhgwJHZIS+vXrFzoBwmvWrNlbb73lz4Gi8vwWiunRRx/NyMgIXQHRkpaWduKJJ9asWXPatGmhW5Lco48+esABB4SugPD222+///W//leFChVmzJgRuoXYcAJAcZxwwgkHHXRQ6AqIovLly1933XXPPfdc6JBk1rlz54MPPjh0BURFpUqV/vrXvz7//POhQ4gNAwDF8bvf/S50AkTawIEDx40bF7oiafXq1St0AkTOgAED3n777WrVqoUOIQYMABSHrxzBXvXu3dsMUEo8dBX2qFu3bh988IE3hbFXHgNKcWzcuLFWrVqhKyAGxo8ff+6554auSDbff/99nTp1QldARC1fvvyss85atGhR6BCiywkARXbkkUfa/UOCnAOUuDZt2tj9w69o2rTptGnTOnbsGDqE6DIAUGSdO3cOnQBx0rt37xdffDF0RfKwBMFeNWzYcPLkyd26dQsdQkQZACiyFi1ahE6AmDn//POfeOKJ0BVJolWrVqETIAbq1q07duzY7t27hw4higwAFJmn70Ex/M///M9dd90VuiIZeAYxJKhWrVrjx4/v2bNn6BAixwBAkTVo0CB0AsTSn/70p+uuuy50Rex5wgkkbv/99//73//uXiD+iwGAIqtZs2boBIil9PT0++67r2/fvqFD4s0SBEVSq1atsWPHmgH4TwYAiuyAAw4InQBxVbly5aeeeur4448PHRJjGRkZoRMgZmrVqvXKK694LhD/ZgCgyKpUqRI6AWLsgAMOGD16dFZWVuiQuKpatWroBIif2rVrT5o0qXXr1qFDiAQvAqPICgoKypf3mwP75JNPPjn22GN3794dOiR+LEFQbN988023bt1WrFgROoTAnABQZD56Yd8dddRRXg5QPJYgKLZDDjlk0qRJ7qPDAAAQRp8+fYYMGRK6AkgtRxxxxJQpU9LS7ABTmusoFFlhYWHoBEgSubm5vXv3/uc//xk6JE4sQbDvxo4d64lkqcz8BxBMpUqVnnzyyczMzNAhQGrp06fP3XffHbqCYJwAUGQuv0HJmjVrlgeDJs4SBCUiPz//iiuuePbZZ0OHEIATAIDAOnfuPGzYsNAVQGpJT09/8MEHXX1ITU4AKDKX36DE5ebmnn322VOmTAkdEgOWIChBK1eu7NSp07p160KHUKYMABSZT18oDatXr27btu3mzZtDh0SdJQhK1nvvvXfiiSeGrqBMuQUIIBIaNWr0zDPPhK4AUs4JJ5zwyCOPhK6gTBkAAKLirLPOGjRoUOgKIOVcddVV/fv3D11B2XELEEXm/B1Kz/r169u2bet+3F9hCYLS8P3333fs2HHFihWhQygLTgAAIqRevXrDhw8PXQGknAMPPPCFF14IXUEZMQAAREuvXr28oRMoe506dfJI4hThFiCKzPk7lLZvvvnmN7/5TU5OTuiQKLIEQenZtWvXGWec8dZbb4UOoXQ5AQCInEMOOcSNQEDZq1y58vDhwytWrBg6hNJlAACIogsvvNCTuYGy16xZM48kTnpuAaLInL9D2Zg9e/Zxxx0XuiJyLEFQ2vLz888+++x//OMfoUMoLU4AACKqU6dOl19+eegKIOWkp6c/9NBDVapUCR1CaTEAAETX4MGD3YwLlL1DDz308ccfD11BaTEAAERXVlbWfffdF7oCSEUXXnhh165dQ1dQKnwHgCJzAy6Upa1btx511FFez/lvliAoMwsWLDjqqKNCV1DynAAARFpGRsadd94ZugJIRUceeeStt94auoKS5wSAInP5DcrYrl272rdvv3jx4tAhkWAJgrK0adOm1q1bZ2dnhw6hJDkBAIi6ypUr//Wvfw1dAaSiWrVqPfTQQ6ErKGFOACgyl9+g7O3evbtTp07z588PHRKeJQjKWF5eXvfu3d95553QIZQYJwAAMVCxYsVbbrkldAWQiipUqDB06NDQFZQkAwBAPPTo0aNjx46hK4BU1LFjx4suuih0BSXGAAAQDxUqVLjhhhtCVwAp6k9/+lNamn1jkvAfEiA2Tj/99MaNG4euAFLR4Ycf/sc//jF0BSXDAAAQG/vtt9+f//zn0BVAirrqqqvS09NDV1ACDAAAcXLOOedUq1YtdAWQijIzMz2SODkYAADipFatWn/6059CVwAp6rLLLnMIkAQMAAAxc/7554dOAFJUvXr1PJI4CRgAAGKmcePG/fv3D10BpKhLL73UIUDcGQAA4ufCCy8MnQCkqIMOOsjTCOKufOgA4qewsDB0AqS6nJyc5s2br1q1KnRIAJYgCG758uXNmjULXUHxOQEAiJ8qVapce+21oSuAFNW0adNBgwaFrqD4nABQZC6/QRSsWLGiadOmoSsCsARBFMyfP79Dhw6hKygmJwAAsdSkSZOzzz47dAWQotq1a3fyySeHrqCYDAAAcdWnT5/QCUCKKl++/DXXXBO6gmJyCxBF5vwdImLDhg0NGjQoKCgIHVKmLEEQEan8NIK4cwIAEFd169b1UjAgFE8jiC8DAECMnXnmmaETgNTVq1ev0AkUhwEAIMa6dOmSlmYlB8Jo3Lhx3759Q1dQZL4DQJG5ATdxV1xxReiE4ihfvnyNGjXKlStXq1atmjVr1qlTp0GDBvXr1z/ooIMqVKgQuo7/NmjQoJEjR4auKDuWoMQNHjw4dEKRZWRkVKxYsWbNmjVq1KhTp069evXq1atXs2bN0F38oqlTp/bo0SN0BUVjAKDIfPomrnz5pPoTq1Klysknn3zccccdddRRbdu2rVevXugiypUrV27ChAm9e/cOXVF2LEGJS5olqGnTpl27dm3Tpk2bNm1atmz58xUKImLHjh2ZmZmbN28OHUIRJMnSQFny6Zu4pPn03aOuXbueddZZJ5988uGHH57c/0sjbvXq1VlZWaEryo4lKHHJ+ofZo0ePnj17du7cuWXLlsn6vzFehgwZcuedd4auoAj82VBkPn0TlyKfTJ06dbryyitPO+20OnXqhG5JUZ06dZozZ07oijJiCUpc0i9Bbdq0GTBgwBlnnNGkSZPQLSnNW4FjJ8mXBkqDT9/EJf2n73+qUqXKTTfdNGDAAJ/EZW/YsGE33XRT6IoyYglKXOosQWedddagQYO6detWsWLF0C2pKD8/v1WrVsuWLQsdQqI8OwIoGTk5OXfeeWfTpk1///vfr1y5MnROauncuXPoBAhp0qRJp59++lFHHTV69OgdO3aEzkk56enpF198cegKiiBVrg1Qglx+S1zqXH77L2lpaXfeeeeVV17p2R1lY8eOHbVr187JyQkdUhYsQYlLzSWoadOmf/3rX3v37l25cuXQLSlk0aJFRxxxROgKEuUEACh5BQUFN9988xFHHPHaa6/ZrpWBqlWr9unTJ3QFRMLy5csvuOCCLl26zJgxI3RLCmnVqlXz5s1DV5AoAwBQWtauXXvGGWdcfPHFGzZsCN2S/I477rjQCRAhH3300UknnXTxxRdnZ2eHbkkJ5cuXHzBgQOgKEmUAAErX888/f+yxx86cOTN0SJJr06ZN6ASInOeee65Vq1Zjx451FFkGTj755NAJJCoV7w5kH1lGE5eaN+D+kieeeOLyyy9PT08PHZKctm3blpGREbqiLFiCEmcJ+rdLLrnk3nvv9ajiUpWbm9uwYcONGzeGDmHvnAAAZeSqq666/vrrd+7cGTokOe2///6nnXZa6AqIqJEjRx5//PELFy4MHZLMKlWq1K9fv9AVJMQAAJSdxx9//IILLti6dWvokOR00kknhU6A6Fq2bNnRRx89bty40CHJrGvXrqETSIgBAChTEydOPPvsszdt2hQ6JAkdddRRoRMg0nbv3t2nT5/7778/Pz8/dEtyateuXegEEmIAAMra9OnT+/Tps2XLltAhycY7mCERf/rTn/7whz/k5uaGDklCjRo1Ovroo0NXsHcGACCA6dOnX3TRRV7YWbIyMzMrVqwYugJi4LHHHrvuuut8Jak0nHXWWaET2DsDABDG5MmTb7jhBgfxJahixYqnnnpq6AqIh6eeeurqq692DlDiOnXqFDqBvTMAAME89dRTjzzySOiKpOJrAJC4UaNGXX/99S5DlCzvJIkFAwAQ0o033jh9+vTQFcmjRYsWoRMgTv72t78NHTo0dEVSycjI8Cyg6DMAAIH179//22+/DV2RJHwPGIrqtttue+6550JXJJUTTjghdAJ7YQAAAtuwYYNT+JLSqFGj0AkQP4MGDZo9e3boiuTRvn370AnshQEACG/8+PFjx44NXZEM6tSpk56eHroCYqagoKBPnz6OIkvKb37zm9AJ7IUBAIiEq6+++rvvvgtdEXvp6ekewg3FsHbt2muvvTYvLy90SDLIyspq2LBh6Ap+jQEAiITNmzfff//9oSuSQcuWLUMnQCxNmjTp6aefDl2RDNLS0s4444zQFfwaAwAQFY888sinn34auiL2GjduHDoB4uqaa65ZtmxZ6Ipk4JHEEWcAACLkrrvuCp0QewcffHDoBIirgoKCG264wY1A+65Zs2ahE/g1BgAgQsaPHz9v3rzQFfF20EEHhU6AGJs6der48eNDV8ReVlZW6AR+jQEAiJannnoqdEK81axZM3QCxNuNN964efPm0BXxdvDBB1esWDF0Bb/IAABEy6hRo1asWBG6IsYyMjJCJ0C8ZWdnjxgxInRFvFWoUMH7gKPMAABEzsSJE0MnxNj+++8fOgFi7/bbb8/Ozg5dEW/t2rULncAvMgAAkfPoo4/m5OSErogrJwCw73JycjwSdB95HViUGQCAyMnOzp47d27oiriqUqWKGQD23X333ff999+HrogxjySOMgMAEEWvv/566IQYO+yww0InQOzl5OSMGzcudEWM1a1bN3QCv8gAAETRc889l5ubG7oirjwJFErEgw8+aCEqNgNAlBkAgCjauHHj4sWLQ1fElafvQYlYuXLl+++/H7oirg444IDatWuHrmDPDABARH300UehE+LqgAMOCJ0ASWL06NGhE2KsQ4cOoRPYMwMAEFEffPBB6IS4qlChQugESBIvvvjixo0bQ1fEVbNmzUInsGcGACCiXn311by8vNAVsVSpUqXQCZAkCgoK3AVUbIccckjoBPbMAABE1Pbt29esWRO6IpaqVasWOgGSx9SpU0MnxJUHEkSWAQCIrpUrV4ZOAFLd6NGjd+zYEboilmrUqBE6gT0zAADRtXr16tAJsbT//vuHToDkkZOTs2jRotAVseSlhJFlAACia+3ataETYmnbtm2hEyCpLFiwIHRCLBkAIssAAESX7wAAUTBr1qzQCbFkAIgsAwAQXdnZ2aETAMq99tpr+fn5oSvip3r16qET2DMDABBd3333XeiEWNq1a1foBEgqW7dudT2iGHwfKbIMAEB0/fjjj6ETYskTS6DEffPNN6ET4qdChQr169cPXcEeGACA6Nq8eXPohFjyc4MS9+2334ZOiCWvAogmAwAQXdu3bw+dEEs5OTmhEyDZeCZB8fgaQDQZAIDoyszMDJ0QS+vWrQudAMnGAEAyMQAA0VWlSpXQCbG0YsWK0AmQbL7//vvQCbFUp06d0AnsgQEAiC43jxbD7t27N27cGLoCko0XkxdP+fLlQyewBwYAILo8PqIYfvrpp9AJkIQ2bNgQOiGWKleuHDqBPTAAANF18MEHh06In02bNoVOAPjffAk4mgwAQHQZAIphy5YtoRMgCX399dehE6DEGACA6PIUoGL44YcfQidAEmrYsGHoBCgxBgAguho3bhw6IX7cqQylwUPJimfXrl2hE9gDAwAQUWlpaVlZWaEr4sezSqA0HHDAAaETYsljCaLJAABE1GmnnVa1atXQFfHjTmUoDZ5KTDIxAAARdcIJJ4ROiKXPP/88dAIkoQMPPDB0Qixt3rw5dAJ7YAAAIuqYY44JnRA/u3bt+vDDD0NXQBLyTILiycvLC53AHhgAgChKT08/8sgjQ1fET3Z2dkFBQegKSEKNGjUKnRBLu3fvDp3AHhgAgCjq379/RkZG6Ir48Q1gKCW+A1A869evD53AHhgAgCg688wzQyfEkm8AQynxXsJiKCgoWLFiRegK9sAAAEROenq6bwAXz+LFi0MnQHI65JBDQifEj2eARpYBAIica665platWqErYmnWrFmhEyAJHX/88Z5KXAwGgMgyAACR07dv39AJsbR169aPPvoodAUkoa5du4ZOiKVt27aFTmDPDABAtHTt2vXoo48OXRFLy5cvD50Ayaldu3ahE2LJABBZBgAgWq6//vry5cuHroilRYsWhU6A5OSpxMXz448/hk5gzwwAQIS0b9++e/fuoSviat68eaETIAl16tTJW8CKZ8uWLaET2DMDABAhd955Z8WKFUNXxFJhYeGUKVNCV0ASOuecc0InxNW6detCJ7BnBgAgKnr06HHqqaeGroirr7/+etWqVaErIAl169YtdEJcrV27NnQCe2YAAKLirrvuSkuzKBXTJ598EjoBklCbNm1atWoVuiKuPJkgsnzWApFw2223tWnTJnRFjH3wwQehEyAJXXbZZR5LUGwLFy4MncCe+Z2myAoLC0MnxIaPjQS1bt36gw8+2H///UOHxFVeXt7BBx+cIrfbWoISZwnad6tXrz744INDV8RSTk7OfvvtF7qCPXMCAIQ3cuRIu/99sWjRohTZ/UNZGjBggN1/sW3YsCF0Ar/IAAAE9vjjj3fo0CF0Rby5/wdKw2WXXRY6IcbWr18fOoFfZAAAQrrwwguvvPLK0BWxN3HixNAJkGy6devWsWPH0BUxtmbNmtAJ/CIDABBM586dH3300fT09NAh8fbNN9/MmDEjdAUkm8GDB/sSxb748ssvQyfwiwwAQBhNmjQZM2ZMjRo1QofE3vvvvx86AZJNt27dTjzxxNAV8fbpp5+GTuAXGQCAAOrXr//666/7dl2JGDt2bOgESDa33Xab15Lsi8LCQieTUeaXGyhr9evXf/vttw8//PDQIclg1apVU6dODV0BSeXiiy8+7rjjQlfE2/r16z2aLMoMAECZatKkybvvvuvNmiXljTfeCJ0ASaVKlSq33HJL6IrYW7VqVegEfo0BACg7nTp1euedd5o3bx46JEkUFBQ8++yzoSsgqTz66KONGzcOXRF7K1asCJ3ArzEAAGWkf//+U6ZMadSoUeiQ5LFw4cKPP/44dAUkj+7duw8cODB0RTJYsmRJ6AR+jQEAKAsPPfTQ888/X7NmzdAhSWXcuHGhEyB5VKlS5YknnqhUqVLokGTw5ptvhk7g13jALUVWWFgYOiE2PEO6XLlyzZs3HzVqlPfplLhNmzbVr19/9+7doUPKmiUocZagInnllVf69OkTuiIZbNiwoV69eqEr+DVOAIBS9Je//GXOnDl2/6VhypQpKbj7h1Jy44032v2XlC+++CJ0AntRIXQAkJy6det29913d+jQIXRIcsrLy3vkkUdCV0CSOO20026//fbQFclj4cKFoRPYCwMAUMJatWo1dOjQnj17pqenh25JWu++++6CBQtCV0AyaNWq1fPPP1+1atXQIclj1qxZoRPYC3cHUmRuwE1cqt2A27FjxxtuuKFnz56VK1cO3ZLkevbs+frrr4euCMMSlLhUW4KKoX79+u+9995hhx0WOiR55ObmHnDAATk5OaFD+DVOAIB9lZaWdumll5533nnHHXecq/5l4KOPPkrZ3T+UoIyMjGnTptn9l6wvvvjC7j/6DABA8Z111lnnnHPOSSedVL9+/dAtKWTYsGGhEyD2MjIypk+f3qZNm9AhyWbu3LmhE9g7AwBQNCeccELXrl07duzYrl27WrVqhc5JOZ988sn48eNDV0C81a9ff9q0aXb/pWHatGmhE9g7AwDw39LS0g499NC0tLSDDz64fv36DRs2zMrKyszMPPTQQ7OysjIyMkIHprSHH344dALEW+vWrSdMmODOn9KwZcuWyZMnh65g7wwAkHJatGjRpUuXVq1aNWrUqH79+nXq1Nl///2rVq3qIRjRN2fOnNGjR4eugBjr3r37888/767FUrJgwYKCgoLQFeydAQBSQps2bc4555yOHTu2bNnSJ19M5efn33bbbaErIMZuuOGGO+64w8WO0jNz5szQCSTEAADJrH379gMHDjz11FObNm3qgYBx98Ybb7z11luhKyCWqlWrNmrUqN69e1sJS09hYeGECRNCV5AQfwYUmYdwJy7UJ01aWtp1113Xr1+/du3apaWlBWmgZG3fvv3YY49dvHhx6JDwLEGJs9n92amnnjp8+PBmzZqFDklyX331lW9WxIUTAEgqVapUue222y644IKGDRuGbqEkPf3003b/UFQVK1YcPnz4wIEDK1WqFLol+bn/J0ZcG6DIXH5LXFlefktLS7vjjjsGDRpUr169Mvt/Stn4+uuvDzvssPz8/NAhkWAJSlyKnwBcdNFFQ4YMOeSQQ0KHpIrf/va3ngEaF04AIBkMGDDg1ltvbdKkSegQSl5BQcEtt9xi9w+J69at25AhQ44//vjQISlk7dq1dv8xYgCAeGvcuPHjjz9++umnp/ilviQ2ceLEl19+OXQFxMPJJ5984403nnLKKenp6aFbUss777wTOoEiMABAjF100UX33ntv3bp1Q4dQWrKzs6+88srQFRADAwYMuOSSS4477jhPPghizJgxoRMoAgMAxFJaWtqzzz47cOBAF/6TWEFBweDBgzdu3Bg6BKKrVatWgwYNOvPMM7OyskK3pK6VK1e++eaboSsoAgMAxE9mZubEiROPPvro0CGUrjFjxrzwwguhKyCK2rdv37dv35NOOqlNmzbu9gnO7j92XDukyDyCI3GlcXn+6KOPHjt2rOdaJL0vvviidevWu3fvDh0SOZagxCXZCeExxxzTtWvXY445pm3btq73R0d+fn6HDh0WLFgQOoQicAIAcdKtW7dXXnmldu3aoUMoXdu3b7/00kvt/klBDRs2rFevXu3atTMzMzMzMxs0aJCVldWoUaNGjRodcMABoevYg08++cTuP3YMABAbPXr0ePHFF2vWrBk6hNJVWFh4yy23eKUOyapFixZdunRp1apVo0aN6tevX6dOnf33379q1apVq1YNnUZxTJo0KXQCRZZUh4OUDefviSvB8/du3bqNHTu2Vq1aJfUPElmjR4++4IILQldElyUocdG5BahNmzbnnHNOx44dW7ZsWb9+/dA5lJgtW7Y0aNAgJycndAhF4wQAYqBNmzYvvfSS3X8qmD9//sCBA0NXQMlo3779RRdddMoppzRt2jQ60wgl6I033rD7jyMDAERd3bp1x48f72H/qeCbb77p1auXl/4Sd2lpadddd12/fv3atWvnqfxJrLCw8MknnwxdQXEYACDqxo4d26xZs9AVlLotW7b0799/7dq1oUOg+KpUqXLbbbddcMEFDRs2DN1CqZs3b55vK8WUAQAi7eGHHz7xxBNDV1Dqdu7cefnll8+ePTt0CBRTWlraHXfcMWjQoHr16oVuoYw899xzoRMoJjfkUWS+gZe4fbzn9Ywzzpg4cWKFCgb1JLd79+5rr732qaeeCh0SD5agxJXZbfcDBgy49dZbmzRpUjb/74iCNWvWNGrUKHQFxWRjARFVrVq1xx57zO4/6RUUFNx88812/8RU48aNH3/88dNPP913fFPNyy+/HDqB4vPVHIioxx9/3Ksuk15hYeFdd931wAMPhA6B4rjooovmzp3bo0cPu/9Us3nz5rvuuit0BcXn4iJEUdeuXT0JPukVFhYOHTp0yJAhoUOgyNLS0p599tmBAwfa+qemCRMmbN26NXQFxefvliJzA27iiv3ROHfu3GOOOaZkY4iU/Pz8IUOG3H333aFD4scSlLhS2p1nZmZOnDjx6KOPLo1/nOjbsWNHixYtVq1aFTqE4nMCAJFzySWX2P0nt9zc3D//+c8PP/xw6BAosqOPPnrs2LGHHHJI6BCCmTRpkt1/3DkBoMhcfktc8S6/LVmypEWLFiUeQ0T89NNPV1111QsvvBA6JK4sQYkr8ROAbt26vfLKK7Vr1y7Zf5YY2bFjR5s2bZYvXx46hH3iBACi5ZJLLrH7T2I//PDDgAEDpk6dGjoEiqxHjx4vvvhizZo1Q4cQ0qRJk+z+k4ATAIrM5bfEFePy24cffujO2mT1+eef9+vX79NPPw0dEm+WoMSV4AlAt27dxo4dW6tWrZL6B4mjn3766cgjjzQAJAEnABAhXbt27dChQ+gKSsX06dPPPvtsz80gjtq0afPSSy/Z/TN27Fi7/+TgPQAQIZdddpln6iWfvLy8xx577JRTTrH7J47q1q07fvz4unXrhg4hsC1btvzlL38JXUHJcAIAUZGenn7qqaeGrqCEff/999dee+0rr7wSOgSKaezYsc2aNQtdQXjPPPPMhg0bQldQMgwAEBUDBw707bokM3fu3AEDBnz55ZehQ6CYHnnkkRNPPDF0BeGtXr168ODBoSsoMW4Bgqg488wzQydQYnJycu67776OHTva/RNfZ5xxxlVXXRW6gkgYNmxYfn5+6ApKjLuNKTKP4EhckW7oX79+vbtsk8PSpUv/53/+57333gsdkpwsQYnbl+8UVatWbcmSJVlZWSXYQ0x99NFHXk+ZZJwAQCR0797d7j8J7Ny58+GHHz7iiCPs/om74cOH2/1Trly5vLw8N/8kH98BgEjw9d8kMHv27Ouuu27+/PmhQ2Bfde3a9fzzzw9dQSSMHz/+nXfeCV1BCXMCAJHQtm3b0AkUX3Z29vXXX3/cccfZ/ZMc7r777goVXCKk3IYNG37/+9+HrqDk+fOGSDj88MNDJ1Ac27dvf/7552+88cacnJzQLVAyLrnkEjd887OhQ4d69GdS8iVgisw38BKX4DfwMjIyNm/enJbmRC5Odu/ePWXKlL/85S/Lli0L3ZJaLEGJK96XgJcsWdKiRYsSjyF2Zs6c2aVLl9AVlAonABDeKaecYvcfI/n5+dOmTbvtttvc8EPyueSSS+z+KVeu3LZt2zwENokZACA8b9mMi927d7/99tv33HPPrFmzQrdAqbjssstCJxAJw4YNW7RoUegKSosBAMJr1KhR6AT2Yvv27a+99trdd9/tE5Ek1rVr1w4dOoSuILwPP/zwjjvuCF1BKTIAQHi1atUKncAvWrt27aRJkx544IG1a9eGboHSddlll+3Lu8NIDlu3br300ktDV1C6DAAQXo0aNUIn8N8KCgrmzJnz4osvjhgxInQLlIX09HQvJKFcuXJ33nmno86k53uHEF716tVDJ/Df0tLSateufeihh2ZmZoZugbIwcODAmjVrhq4gsDfffHPYsGGhKyh1BgAIr2LFiqET2IPmzZv/8Y9/XL58+dtvvz1gwIDQOVC6zjzzzNAJBPbtt98OHDgwdAVlwQAA8GsqV67crVu3559/fvny5UOHDs3IyAhdBKXi6KOPDp1ASHl5eVdfffW6detCh1AWDAAACWnSpMnNN9+8Zs2a0aNHt2rVKnQOlKTu3bvXrVs3dAUhPf7445MnTw5dQRkxAEB4BQUFoRNIVEZGRv/+/efPn//qq6+2b98+dA6UDF//TXEzZsz4wx/+ELqCsmMAgPB++umn0AkUTeXKlc8888w5c+a88cYbnTp1Cp0D+6pt27ahEwhmzZo1/fv3D11BmTIAQHg//vhj6ASKo0KFCqeddtqMGTMmTpzYokWL0DlQfIcffnjoBMLYsWPHoEGDsrOzQ4dQpgwAEJ4BINYqVqx41llnzZs377nnnqtfv37oHCiyjIyMBg0ahK4ggMLCwr/85S9vvfVW6BDKmgEAwlu9enXoBPZV1apVBw4cuGjRoltuuSV0CxTNKaeckpZmP5CKRo4c+eijj4auIAB/8BDeihUrQidQMurUqXPnnXd+9tlnPXv2DN0CiWrWrFnoBAKYPn36pZdeGrqCMAwAEN4HH3wQOoGS1Lp168mTJ7/yyiu1a9cO3QJ716hRo9AJlLUlS5acffbZoSsIxgAA4a1cufKHH34IXUFJSk9P79Onz2effea1mkRfrVq1QidQprKzs88666ytW7eGDiEYAwBEgruAklLDhg1HjRo1efJkRwFEWY0aNUInUHY2b97ct2/fL7/8MnQIIRkAIBI+++yz0AmUivLly//ud79buHDhWWedFboF9qx69eqhEygjO3fuvPTSS2fOnBk6hMAMABAJ77//fugESlFmZubYsWOfeOIJz1ohgipWrBg6gbKQm5t7/fXXT5w4MXQI4fkogkgYO3as9wEntwoVKvzP//zP3LlzmzZtGroFSDn5+fmDBw8eMWJE6BAiwQAAkbB79+4FCxaErqDUdejQYdasWW4HAspSYWHh7bff/tBDD4UOISoMABAVb775ZugEykK9evVefvnlO+64I3QI/G8FBQWhEyhFhYWFQ4cOvfPOO0OHECEGAIiKZ555ZteuXaErKAuVKlW69dZbx4wZ4ysBRIH7D5PYz7v/IUOGhA4hWnz2QFRs2LDBkxlSSr9+/WbMmFGzZs3QIaS6H3/8MXQCpcLun19iAIAIeemll0InUKaOP/742bNnH3bYYaFDSGkGgKSUn58/ZMgQu3/2yAAAEfL8889//fXXoSsoU82bN58+fXr79u1Dh5C6Vq9eHTqBEpabmzt48OChQ4eGDiGiDAAQLQ4BUtDBBx88ZcqUrl27hg4hRXkTeZLZuXPnVVdd9cADD4QOIbrKhw4gfgoLC0MnxEb58kX+E6tSpco333xTr1690ughyjZt2tSnT5/p06eHDok6S1DiElyCGjdu7OwxaWzZsuWSSy6ZNGlS6BAizQkAREtOTs7IkSNDVxBArVq1xo4d261bt9AhpJyVK1f+8MMPoSsoAdnZ2b169bL7Z6+cAFBkLr8lrhgnAOXKlUtPT1+xYkVWVlaJ9xB9zgH2yhKUuMSXoLlz5x5zzDGlGkNpW7p0aa9evb788svQIcSAEwCInPz8/HvuuSd0BWHUqlXr5ZdfthWjjH322WehE9gnM2bM6Nixo90/CTIAQBSNGDFixowZoSsIo06dOuPGjWvSpEnoEFLI+++/HzqBYiosLHz22WdPOumkrVu3hm4hNtwCRJE5f09c8W4B+lmrVq1mz569//77l2APMbJ06dLOnTtv3rw5dEjkWIISl/gSVLFixU2bNlWvXr1UeyhxO3fuvOWWWx566KHQIcSMEwCIqMWLF991112hKwimRYsWU6ZMSUuzSlMWdu/evWDBgtAVFM3atWt79epl908x+GiB6LrvvvvefPPN0BUE06lTpxdeeCF0BanCahMvs2bN6tixo/9qFI9bgCgy5++J25dbgH5Wu3btDz/80O3gqWzIkCF33nln6IoIsQQlrkhLUN26dVevXl25cuXS66FE5OXl/e1vf7vuuutChxBjTgAg0jZu3HjRRRf5alcqu+WWW84+++zQFSS/DRs2zJw5M3QFe7Fu3bp+/frZ/bOPDAAQdTNnzrzyyitzc3NDhxBGpUqVnnjiCadAlIGXXnopdAK/5l//+texxx47YcKE0CHEnluAKDLn74nb91uA/u2GG26477770tPTS+ofJF7mzJnTqVOn0BWRYAlKXDGWoBUrVhx66KGlEcO+2L59+7333jt06NDQISQJJwAQDw8++OBf//rX/Pz80CGE0bFjx+HDh4euIPk5BIigTz755OSTT7b7pwQ5AaDIXH5LXAmeAPzspptuuueee5wDpKa8vLzevXtPnjw5dEhglqDEFWMJqlKlyjfffFOvXr3S6KGoduzY8dhjjw0ePDh0CMnGCQDEyQMPPHD11Vfv3LkzdAgBVKhQ4ZFHHqlZs2boEJJZTk7OqFGjQldQrly5cvPmzevatavdP6XBCQBF5vJb4kr8BOBnPXv2HDVq1IEHHlga/zgRN378+HPPPTd0RUiWoMQVbwlKT09fsWJFVlZWifeQoC1btgwbNsy7ICk9TgAgfqZMmdKlS5eFCxeGDiGAc845p3///qErSGb5+fn33HNP6IoUVVBQ8Nprr7Vv397un1LlBIAic/ktcaV0AvCz9PT0p59+esCAAb4SkGpWr17dokWL7du3hw4JwxKUuH1Zgt59990TTzyx5FrYu2XLlt1yyy0TJ04MHULycwIAcZWfn3/JJZecd95533zzTegWylSjRo0ee+yx0BUkuWuuuWbbtm2hK1LFxo0bb7/99pYtW9r9UzYMABBv48aN+81vfvPss8/m5OSEbqHsXHjhhV27dg1dQTJbvHixu1DKwK5du55//vkWLVrcdtttBQUFoXNIFW4BosicvyeuVG8B+i/t2rW79957u3btmpZmsE8JH3300THHHBO6IgBLUOL2fQmaNm1a9+7dSySG/5Kfnz916tRbb731008/Dd1CyrFRgCTx8ccfn3LKKd27d3/nnXdcRkoFRx999NVXXx26giTXv3//r7/+OnRFsikoKHjnnXe6det2xhln2P0ThBMAiszlt8SV5QnAf+rUqdMf/vCH7t27V69ePUgAZWPt2rXNmjVLtbu/LEGJK5El6Pjjj58yZUpGRsa+/1MUFha+9957d9111/Tp00O3kNIMABSZT9/EhRoAfpaRkXHdddedfvrp7dq1q1ixYsASSs9DDz10ww03hK4oU5agxJXUEnTeeec999xzlSpVKpF/LTXl5+e/++67Dz/88NSpU0O3gAGAovPpm7iwA8C/ZWZmDhw4sEuXLm3btvX6sCSzefPm5s2bb9iwIXRI2bEEJa4El6Abbrjhvvvu89DhYti1a9e0adOGDRs2a9as0C3wv0Vid0K8+PRNXEQGgP/Uvn377t27N2/evFmzZg0aNKhbt26VKlVCR7FPnnnmmcsuuyx0RdmxBCWuZJegm2+++fbbbzcDJG7Tpk2vvvrq/fff/+WXX4Zugf9L5HYnRJ9P38RFcAD4f2VlZTVu3LhOnTq1atX6leBKlSpVq1atYsWKNWvWrFmzZv369evXr9+wYcMDDzwwFv8zk9hPP/30m9/8Zu3ataFDyoglKHEl/rd500033XPPPWaAvfr888/Hjh173333pdpXdIgLH9sUmU/fxKXCzjgzM/Pkk08+9thj27Zte8QRR1StWjV0USoaNWrUJZdcErqijFiCElcaS9AVV1zx0EMP7bfffiX+LyeBnTt3vvvuu08//fQ//vGP0C3wa5J/d0KJ8+mbuFQYAP5TxYoVzzvvvB49epxwwgl169YNnZNCfvzxx8MOOyxFvglgCUpcKS1BPXv2HDVqlC8U/aelS5e++uqrw4cPX7duXegW2LvU2p1QInz6Ji7VBoD/NHDgwP79+3fp0sWTQ8rG8OHDr7nmmtAVZcESlLjSW4KaN2/+8ssvt23btpT+/bj47rvv3n777b///e/vvPNO6BYogtTdnVBsPn0Tl8oDwM+aN29+0003nXXWWTVq1AjdkuS+//77gw46aPfu3aFDSp0lKHGlugSlp6c//fTTAwYMSMGvBPzwww8zZ86cMGHCmDFjQrcAlIlCEhb6v1VU1K1b929/+9uPP/4Y+j9Ikhs8eHDo/9RlIfSPOU7K4D/Hueeeu3LlytD/Q8vI2rVrx40b17dv3zL4wUKpSvXLkxRDoX1twpwA/KfGjRvff//9Z555ZoUKFUK3JKelS5e2bNkydEWpswQlrmyWoCpVqjz++OPnn39+Uj5TOD8///PPP581a9arr7761ltvhc6BkmF3QpH59E2cAeD/ddpppz344IMtWrQIHZKczjnnnIkTJ4auKF2WoMSV5RLUrl27e++9t2vXrmlpaWX2/7T0rF27dsGCBe+9997LL7+cnZ0dOgdKmN0JRebTN3EGgD1KS0t79NFHL7/88ooVK4ZuSTZTp07t0aNH6IrSZQlKXNkvQd26dRs8ePCJJ54YuzGgsLBwzZo1S5cunTt37pQpUz7++OPQRVCK7E4oMp++iTMA/IqePXs++eSTmZmZoUOSyo4dO5o0aZLcDyK0BCUu1BLUqVOnP/zhD927d69evXqQgAT9+OOPy5cvX7p06bx586ZMmbJy5crQRVBG7E4oMp++iTMA/LrMzMyxY8d26tQpdEhSufvuu2+++ebQFaXIEpS4sEtQRkbGddddd/rpp7dr1y4Kx335+fnffffdypUrf970z5w588MPPwwdBWHYnVBkPn0TZwDYq7S0tBdffPG8884LHZI8kv6rwJagxEVkCcrKyrrgggu6dOnStm3bsnl9WH5+/g8//LB+/frvvvtu1apVy5cvX7Ro0cyZM7dv314G/98h+iKxNBAvPn0TF5FP3+i77777brrpJj+uElFYWHjcccfNmTMndEhpsQQlLoJ/U0cfffQpp5zSvHnzZs2aNWjQoG7dusV4dtCuXbt27ty5Y8eOrVu3/vjjj1u2bNmyZcv69evXrVu3Zs2ar776av78+fn5+aXRD8khcksD0efTN3ER/PSNrNtvv/3mm29OwTcKlYann3768ssvD11RWixBiYvFEpSVldW4ceM6derUqlVrj8GbN28uKCj44Ycf8vLy1q1bt2rVqlR44R2UqhgsDUSNT9/ExeLTNzruuOOOW265xQ9t361cufLQQw8NXVFaLEGJ89cE7FHMntIFJLEhQ4Y8/PDDoSuSQePGjU8++eTQFQBElAEAiJAbbrhh7NixoSuSQe/evUMnABBRDgcpMufviXP+XgxpaWmzZ88+5phjQofE2+eff56sr1u2BCXOEgTskRMAIFoKCgr69euX3K+yKgPNmzdv3bp16AoAosgAAETOypUrr7nmmry8vNAhMVa+fHl3AQGwRwYAIIomTJjw3HPPha6It86dO4dOACCK3B1IkbkBN3FuwN0XFStW/Pzzz5s0aRI6JK62bNlSu3btgoKC0CElzBKUOEsQsEdOAICI2r179x//+Mfk27+WmRo1apxxxhmhKwCIHAMAEF2TJk164403QlfE2EknnRQ6AYDIMQAAkfbHP/5xx44doSvi6qijjgqdAEDkGACASFu6dOmkSZNCV8RVy5YtQycAEDkGACDqbr/9docAxVOzZs0TTzwxdAUA0WIAAKJu+fLlr7/+euiKuDr++ONDJwAQLQYAIAYefPDB/Pz80BWx5H3AAPwXAwAQAx9++OH8+fNDV8RS8+bNQycAEC0GACAexo0bFzohlrxJDYD/4h2BFJnXcCbOazhLULVq1davX1+tWrXQIfFz7LHHfvjhh6ErSowlKHGWIGCPnAAA8bB9+/Y5c+aEroilDh06hE4AIEIMAEBsTJs2LXRCLLVo0SJ0AgARYgAAYuPll1/Oy8sLXRE/vgYAwH8yAACxkZ2d/eWXX4auiJ8GDRqETgAgQgwAQJwsWrQodEL8GAAA+E8GACBOvA2gGGrXrp2RkRG6AoCoMAAAcZJMj7MsM+XLl/cgIAD+zQAAxMnMmTN37twZuiJ+DjnkkNAJAESFAQCIme+++y50QvxkZmaGTgAgKgwAQMysX78+dEL8NGzYMHQCAFFhAABi5vvvvw+dED+1atUKnQBAVBgAgJjZsmVL6IT4OeCAA0InABAVBgAgZn788cfQCfFjAADg3wwAQMxs3749dEL8eA8AAP9mAABIfvvtt1/oBACiwgAAkPwMAAD8mwEAIPlVqlQpdAIAUWEAAGKmevXqoRPixwAAwL8ZAICY8UCbYqhcuXLoBACiwgAAxEzNmjVDJ8RPenp66AQAosIAAMTMgQceGDoBAGLMAADETN26dUMnAECMlQ8dQPwUFhaGToiN8uX9iZWwtLS07du3V6lSJXRI/CTNb6MlKHFJ8x8dKFlOAIA4Of744+3+iyEvLy90AgBRYQAA4uTYY48NnRBLubm5oRMAiAoDABAn7dq1C50AAPFmAADipHXr1qETYmnHjh2hEwCICgMAEBtZWVnNmjULXQEA8WYAAGKjb9++XmhVPL4DAMC/GQCA2OjWrVvohLjatm1b6AQAosIAAMRDtWrVOnbsGLoirrZu3Ro6AYCoMAAA8XD11VdXq1YtdEVcbdmyJXQCAFFhAADi4ZxzzgmdEGM//vhj6AQAosIAAMTA8ccff9RRR4WuiDEnAAD8mwEAiIHf//73aWnWq+LbvHlz6AQAosIHKhB1TZs2Pf3000NXxNt3330XOgGAqDAAAFE3dOjQKlWqhK6ItxUrVoROACAqyocOIH4KCwtDJ8RG+fL+xPbVkUceOWfOnMqVK4cOibcGDRqsW7cudEXJsAQlzhIE7JETACDS7r33Xrv/fbR169ak2f0DsO8MAEB09evX75RTTgldEXsbNmwInQBAhBgAgIiqUqXK3Xff7R6GfefyPwD/yQAARNSIESMOOeSQ0BXJIDs7O3QCABFiAACiqH///v379w9dkSQ8AgiA/2QAACKnefPmDz30UHp6euiQJPHJJ5+ETgAgQtxcS5F5Bl/i3L9eDBUrVpw7d+5RRx0VOiRJFBQU1KlTJ5neBGwJSpwlCNgjJwBAtIwfP97uvwRlZ2cn0+4fgH1nAAAiZPjw4b/73e9CVySVlStXhk4AIFoMAEBU3HfffVdddVXoimTz1VdfhU4AIFoMAEAk3HvvvTfddFPoiiS0aNGi0AkAREuF0AEA5R5//PGrr746dEUSKiwsnDp1augKAKLF8wEoMo/gSJxHcOxVenr6+PHje/XqFTokOa1Zs6ZRo0ahK0qYJShxliBgj9wCBATTtGnTuXPn2v2XniVLloROACByDABAGOeee+6sWbPat28fOiSZffzxx6ETAIgc3wEAylrFihWffPLJgQMHetdvaZs+fXroBAAix92BFJkbcBPnBtz/11lnnXX//fc3adIkdEjy27JlS82aNUNXlDxLUOIsQcAeOQEAykiLFi3uv//+3/72t2lpbj4sCwsXLgydAEAUGQCAUpeZmTl06NDevXtXrVo1dEsKee+990InABBFBgCgFLVr1+7GG2/s2bNn9erVQ7eklsLCwnHjxoWuACCKDABAyUtPT7/yyiv79u17zDHHVKhgnQlg+fLlS5cuDV0BQBT5YAZKTM2aNc8777zTTjutc+fONWrUCJ2T0ubOnRs6AYCIMgAA+6R169YnnHBChw4djjjiiBYtWlSqVCl0EeXKlSv3xhtvhE4AIKI8IIwi8wy+xA0ePDh0QkmqUqVKlSpVatasWaNGjbp169arV69hw4YHHHBA6C7+25YtW+rUqZOfnx86pFRYghLnMaDAHlkaKDKfvhBxr7322hlnnBG6orRYghJnAAD2yNO4AZLNhAkTQicAEF2uDVBkLr9BlG3cuLFu3boFBQWhQ0qLJShxTgCAPXICAJBU3n///STe/QOw7wwAAEnF+78A+HUOByky5+8QWevWrWvQoEHoitJlCUqcW4CAPXICAJA8pkyZEjoBgKhzbYAic/kNoik/P79du3affvpp6JDSZQlKnBMAYI+cAAAkiXnz5iX97h+AfWcAAEgSY8aMCZ0AQAw4HKTInL9DBG3YsKFhw4b5+fmhQ0qdJShxbgEC9sgJAEAyeP3111Nh9w/AvnNtgCJz+Q2iZteuXW3btl22bFnokLJgCUqcEwBgj5wAAMTe22+/nSK7fwD2nQEAIN4KCwsffPDB0BUAxIbDQYrM+TtEyuzZs4877rjQFWXHEpQ4twABe+QEACDe/va3v4VOACBOXBugyFx+g+hYunRpy5YtQ1eUKUtQ4pwAAHvkBAAgxh566KHQCQDEjGsDFJnLbxARCxcuPPLII0NXlDVLUOKcAAB75AQAIK48/AeAYnBtgCJz+Q2i4KOPPjrmmGNCVwRgCUqcEwBgj5wAAMTSsGHDQicAEEuuDVBkLr9BcDNmzDjppJNCV4RhCUqcEwBgjywNFJlPXwgrNze3S5cuH374YeiQMCxBiTMAAHvkFiCAmHn55ZdTdvcPwL5zbYAic/kNAvrhhx9atmy5YcOG0CHBWIIS5wQA2CMnAABx8thjj6Xy7h+AfefaAEXm8huEsnTp0pYtW4auCMwSlDgnAMAeOQEAiIe8vLybbropdAUAsWcAAIiHl156aerUqaErAIg9h4MUmfN3KHurV69u0aLF9u3bQ4eEZwlKnFuAgD1yAgAQdYWFhYMHD7b7B6BEuDZAkbn8BmVs8uTJvXr1Cl0RFZagxDkBAPbI0kCR+fSFsrRu3bojjzxy3bp1oUOiwhKUOAMAsEduAQKIroKCghtvvNHuH4ASZAAAiK4xY8a89NJLoSsASCoOByky5+9QNpYtW3bEEUfs3r07dEi0WIIS5xYgYI+cAABE0fbt2y+99FK7fwBKnAEAIIruuuuuWbNmha4AIAk5HKTInL9DaZs4ceI555wTuiKiLEGJcwsQsEeWBorMpy+Uqs8///yoo47KyckJHRJRlqDEGQCAPXILEECEbNq06YILLrD7B6D0GAAAoiIvL++aa675+OOPQ4cAkMwMABSZ83coJffdd9+YMWNCVwCQ5NwdSJFt3769atWqoSsg2YwdO7Zv376hK2LANYjE+Q4AsEdOACiyHTt2hE6AZDN79uzzzjsvdAUAKcEAQJFt3bo1dAIklWXLlvXq1augoCB0CAApwQBAkW3evDl0AiSP7OzsXr16bdiwIXQIAKnCAECRZWdnh06AJLFp06Y+ffosW7YsdAgAKcQAQJF9++23oRMgGWzbtu2CCy6YNWtW6BAAUosBgCJbvHhx6ASIvZ07d1522WVTp04NHRI/+fn5oRPiwQ8K+CUGAIrMBUvYRz/v/l955ZXQIbG0a9eu0Anx4AcF/BIDAEX26aef/vDDD6ErIK5yc3Ovueaa0aNHhw6Jq23btoVOiAc/KOCXGAAojiVLloROgFjauXPnlVdeOXLkyNAhMWZfmyA/KOCXGAAojo8//jh0AsTPz3f+jBo1KnRIvDmBTNDGjRtDJwARZQCgOF599dXQCRAz27Ztu/jii935s+88iThBntgG/BIDAMUxa9asNWvWhK6A2Ni0adO5557rW78lYvXq1aET4sEPCvglBgCK6b333gudAPHw7bff9uzZc9q0aaFDksSiRYtCJ8SDHxTwSwwAFNPLL78cOgFiYNmyZV27dp0zZ07okOTx7rvvhk6IBz8o4JeUDx1AjC1fvrxJkyahKyC6Zs2adeaZZ/ouZolbu3btQQcdFLoi0r799tvMzMzQFUBEOQGg+CZNmhQ6AaJrzJgxJ5xwgt1/afA+8r3yIwJ+hQGA4nvkkUd27NgRugIiZ/fu3XfeeWf//v0LCgpCtySnjz76KHRC1PkRAb/CAEDxZWdnv/baa6ErIFo2btx44YUXDhkyJHRIMhs3blxhYWHoiugqLCwcN25c6AogunwHgH3Svn372bNnV6xYMXQIRMLSpUvPP//8BQsWhA5Jfl9++WWzZs1CV0TUV199ddhhh4WuAKLLCQD7ZP78+VOnTg1dAZEwbty4tm3b2v2XjXfeeSd0QnT54QC/zgDAvrrllltycnJCV0BIP/3000033dSnT5/du3eHbkkVo0ePdhfQL/HCaeDXGQDYV4sXLx4zZkzoCghm6dKlp5566rBhw0KHpJZZs2Z50M0eLVq0aNasWaErgEgzAFACbrzxxg0bNoSugLKWn5///PPPt23b1nu+ghg/fnzohCjyYwH2ypeAKRlXX331448/HroCyk52dvYNN9zwyiuvhA5JXRkZGatWrapRo0bokAjZsmVLVlbW1q1bQ4cAkeYEgJIxfPjwf/3rX6EroCwUFBRMmDDhiCOOsPsPa+vWrV5H+F9effVVu39gr5wAUGIaN248d+7cunXrhg6BUrRy5co///nPHrIeEU2bNv3000+rVq0aOiQSduzY0aZNm+XLl4cOAaLOCQAlZuXKlX/605+8+pRktXv37pEjR7Zo0cLuPzqWL1/uEODfJk2aZPcPJMIJACXsySefvOKKK0JXQAn79NNPb7zxxunTp4cO4b81bdp0wYIF1atXDx0S2E8//XTkkUcaAIBEOAGghF155ZUzZ84MXQElZv369YMHD27btq3dfzQtX7585MiRoSvCGzlypN0/kCAnAJS82rVrf/DBB4cffnjoENgnOTk5L7300k033bR58+bQLfyaKlWqfPHFF40aNQodEszq1asPP/xw72QEEuQEgJK3cePGc845Z926daFDoJgKCwunT5/eqVOnQYMG2f1HX05Ozq233pqyLwYuLCy89dZb7f4BCK9jx47r1q0rhLiZNWtWz549Q/8BUWSTJ08O/bsTxj/+8Y/QP3sgZtwCRCnq1q3b6NGj69WrFzoEEvLRRx8NGzbMi1Rjqm7dugsWLGjYsGHokDKVnZ3drl07J64AREjHjh2//fbb0BfIYC8WLFhw4YUXhv5zYV/17t179+7doX+bys7u3bt79+4d+qcOAP+PFi1afP7556E/KGEP8vLyZsyYcfbZZ4f+K6HEDBs2LPSvVdkZNmxY6J83APyCmjVrzpgxI/RnJfwfO3bsmDBhQseOHUP/cVDy3nzzzdC/X2XhzTffDP2TBoC9eeKJJ/Ly8kJ/aJLqsrOzn3zyycaNG4f+g6C0ZGRkfPbZZ6F/0UrXZ599lpGREfonDQAJ6N+/v0cDEcTu3bvffffdyy67LC3N44+TX5MmTVavXh36l660rF69ukmTJqF/xgCQsMzMzGnTpoX+ACWFrFq16oknnmjevHno333KVOvWrb/77rvQv30lb926da1btw790wWAorv88ss9HYhS9e23344ZM+bMM88M/ctOMMn3NpJ169b54goAMVatWrUnn3xy+/btoT9SSSrr168fN25c3759Q/+CEwlt2rRJmnuBVq9e3aZNm9A/UQDYZ4cddti4ceNycnJCf7YSY3l5eYsWLRoxYsTvfve70L/RRE7Tpk0XL14c+pd0Xy1ZsqRp06ahf5YAUHJatWr197//fevWraE/ZImTdevWTZs27c9//nNWVlboX2EiLSMjY/r06aF/YYtv+vTpnvkDQHKqXbv2HXfcsWTJktCftkRUfn7+8uXLX3311Ztvvrldu3ahf2GJmUcffTR2TyLOy8t79NFHQ//kgGRTPnQA7EHnzp0vuuiiLl26NGnSpHx5v6WpKzc3d82aNStXrvzyyy8//fTTf/7zn+vWrQsdRYz169fv4YcfrlevXuiQhKxfv/76669/+eWXQ4cAycbWikhr3br1ueee26FDh1atWh100EGhcyhFBQUFmzZtWrdu3bp167Kzs7/66qv58+f/61//2r17d+g0kkpmZuazzz7bvXv30CF78eabbw4aNGjt2rWhQ4AkZAAgNrKysk488cRWrVplZWU1bNiwdu3aGRkZ1atXr1ChQnp6euXKlUMH8t927NhRrly53NzcvLy8goKCnTt3FhQUbP3/bd68edOmTZs2bdqwYcPSpUs//PDDnJyc0Mmkissvv3zIkCENGzYMHbIH2dnZd9xxx4gRI0KHAABAEsnIyBgxYsSOHTtC3+f/f+zcuXPEiBG+7wsAAKWlVatW48ePz83NDbv1z83NHT9+fKtWrUL/PAAAIAW0adNmzJgx27ZtK/ut/7Zt28aMGeMNXwAAUNbq169/zz33fPHFF2Wz9f/iiy/uueee+vXrh/7fDQAAqa179+5///vfV61aVRr7/lWrVv3973+P/mOIgCTmKUAAsGcdO3bs1atXhw4dWrZseeCBBxb73/n++++XLFkyf/78SZMmzZkzpwQLAYrBAAAAe9eqVasuXbq0atUqMzOzQYMGtWrVysjIqFq16n777ffz+woLCwt37ty5Y8eOrVu3btq06bvvvlu7du3ixYvff//9xYsXh84HAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAKA3/H0WGFEv6h/PwAAAAAElFTkSuQmCC"
        # gui.theme(THEME)
        gui.theme(THEME)
        self.window = gui.Window(
            "Create Sudokus...",
            self.l_caf(),
            background_color="#000000",
            icon=(self.bpicon_base64),
            finalize=True,
            resizable=True,
        )
        while True:
            event, values = self.window.read()

            if event == "Exit" or event == "Cancel" or event == gui.WIN_CLOSED:
                print("Closing:", self.window.Title)
                self.window.close()
                self.closed = True
                break

            if event == "-EP1-":
                NO1 = values["-EP1-"]
                if NO1 == True:
                    self.window['-EP2-'].update(value=False)
                    self.window['-EP3-'].update(value=False)
                    self.window['-EP4-'].update(value=False)
                    self.window['sec_1'].update(visible=True)
                    self.window['sec_2'].update(visible=False)
                    self.window['sec_3'].update(visible=False)
                    self.window['sec_4'].update(visible=False)
            if event == "-EP2-":
                NO2 = values["-EP2-"]
                if NO2 == True:
                    self.window['-EP1-'].update(value=False)
                    self.window['-EP3-'].update(value=False)
                    self.window['-EP4-'].update(value=False)
                    self.window['sec_1'].update(visible=False)
                    self.window['sec_2'].update(visible=True)
                    self.window['sec_3'].update(visible=False)
                    self.window['sec_4'].update(visible=False)

            if event == "-EP3-":
                NO3 = values["-EP3-"]
                if NO3 == True:
                    self.window['-EP1-'].update(value=False)
                    self.window['-EP2-'].update(value=False)
                    self.window['-EP4-'].update(value=False)
                    self.window['sec_1'].update(visible=False)
                    self.window['sec_2'].update(visible=False)
                    self.window['sec_3'].update(visible=True)
                    self.window['sec_4'].update(visible=False)

            if event == "-EP4-":
                NO4 = values["-EP4-"]
                if NO4 == True:
                    self.window['-EP1-'].update(value=False)
                    self.window['-EP2-'].update(value=False)
                    self.window['-EP3-'].update(value=False)
                    self.window['sec_1'].update(visible=False)
                    self.window['sec_2'].update(visible=False)
                    self.window['sec_3'].update(visible=False)
                    self.window['sec_4'].update(visible=True)

            elif event == "Ok":

                global exportfile
                global exportfile2
                global exportfile3
                global PPTX
                global PNG
                global JPG
                global SVG
                global difficulty
                global noofpuzzles
                global sudokusolutiontitle
                global sudokutitle

                PPTX = False
                PNG = False
                JPG = False
                SVG = False

                difficulty = values["-DIFF-"]
                noofpuzzles = values["-SUDOKUNO-"]
                sudokusolutiontitle = values["-SUDOKUSOLUTIONTITLE-"]
                sudokutitle = values["-SUDOKUTITLE-"]

                PPTX = values["-EP1-"]
                PNG = values["-EP2-"]
                JPG = values["-EP3-"]
                SVG = values["-EP4-"]

                exportfile = values["-EXPORTFILE-"]
                if exportfile == '':
                    exportfile = 'qqwing.exe'

                if PPTX == True:
                    exportfile2 = values["-EXPORTFILE2-"]
                if PNG == True:
                    exportfile2 = values["-EXPORTFILE3-"]
                    exportfile3 = values["-EXPORTFILE6-"]
                if JPG == True:
                    exportfile2 = values["-EXPORTFILE4-"]
                    exportfile3 = values["-EXPORTFILE7-"]
                if SVG == True:
                    exportfile2 = values["-EXPORTFILE5-"]
                    exportfile3 = values["-EXPORTFILE8-"]

                confirm = "true"

                if PNG == True or JPG == True or SVG == True:
                    if exportfile3 == '':
                        gui.Popup(
                            "You need to choose the save location for your file(s)!"
                        )
                        confirm = "false"

                if exportfile2 == "":
                    gui.Popup(
                        "You need to choose the save location for your file(s)!"
                    )
                    confirm = "false"

                try:
                    int(noofpuzzles)
                except:
                    gui.Popup(
                        "Your amount of puzzles to generate must be a number!"
                    )
                    confirm = "false"
                    
                if confirm == "true":
                    if int(noofpuzzles) < 1:
                        gui.Popup(
                        "You must generate at least 1 puzzle!"
                        )
                        confirm = "false"

                if confirm == "true":
                    break

    def l_caf(self):

        powerpointfile = [
            [gui.Push(), gui.Text("Choose where you want to save your PowerPoint file", font=('italic', 12)), gui.Push()],

            [
                gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE2-"),
                gui.FileSaveAs(file_types=(("PPTX files", "*.pptx"),),
                                default_extension = "*.pptx",),
                gui.Push(),
            ],
        ]

        pngfile = [

            [gui.Push(), gui.Text("Choose where you want to save Puzzle PNG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE3-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],

            [gui.Push(), gui.Text("Choose where you want to save Solution PNG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE6-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],
        ]

        jpgfile = [

            [gui.Push(), gui.Text("Choose where you want to save your Puzzle JPG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE4-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],

            [gui.Push(), gui.Text("Choose where you want to save your Solution JPG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE7-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],
        ]

        svgfile = [
            [gui.Push(), gui.Text("Choose where you want to save your Puzzle SVG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE5-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],

            [gui.Push(), gui.Text("Choose where you want to save your Solution SVG files", font=('italic', 12)), gui.Push()],
            [gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE8-"),
                gui.FolderBrowse(),
                gui.Push(),
             ],
        ]

        pdf_column = [
            [
                gui.Push(),
                gui.Text("QQWing", font=("Arial", 18, "bold")),
                gui.Push(),
            ],
            [
                gui.Push(),
                gui.Text(
                    "Please choose the location of the QQWing.exe. If left blank the program assumes it's in the same directory as this one."
                ),
                gui.Push(),
            ],
            [
                gui.Push(),
                gui.In(size=(25, 1), enable_events=True, key="-EXPORTFILE-"),
                gui.FileBrowse(
                    file_types=(
                        ("QQWing", "qqwing.exe"),
                    )
                ),
                gui.Push(),
            ],
            [
                gui.Push(),
                gui.Text("Sudoku Settings", font=("Arial", 18, "bold")),
                gui.Push(),
            ],

            [
                gui.Push(),
                gui.Text(
                    "How many Sudoku puzzles would you like to make?"
                ),
                gui.Push(),
            ],

            [
                gui.Push(),
                gui.Input(key='-SUDOKUNO-', size=(4, 1), default_text=50),
                gui.Push(),

            ],

            [
                gui.Push(),
                gui.Text(
                    "What difficulty level?"
                ),
                gui.Push(),
            ],

            [
                gui.Push(),
                gui.Combo(self.difficultylevel,default_value=self.difficultylevel[0], enable_events=True, readonly=True, size=(8, 1), key='-DIFF-'),
                gui.Push(),

            ],


            [
                gui.Push(),
                gui.Text(
                    "Sudoku Puzzle Title"
                ),
                gui.Input(key='-SUDOKUTITLE-', size=(16, 1), default_text="Sudoku"),
                gui.Text(
                    "Sudoku Solution Title"
                ),
                gui.Input(key='-SUDOKUSOLUTIONTITLE-', size=(16, 1), default_text="Sudoku Solution"),
                gui.Push(),
            ],
            

            
            [
                gui.Push(),
                gui.Text("Export Settings", font=("Arial", 18, "bold")),
                gui.Push(),
            ],
            [
                gui.Push(),
                gui.Checkbox(
                    "PPTX File",
                    key="-EP1-",
                    default=True,
                    enable_events=True,
                ),
                gui.Checkbox(
                    "PNG File",
                    key="-EP2-",
                    default=False,
                    enable_events=True,
                ),
                gui.Checkbox(
                    "JPG File",
                    key="-EP3-",
                    default=False,
                    enable_events=True,
                ),
                gui.Checkbox(
                    "SVG File",
                    key="-EP4-",
                    default=False,
                    enable_events=True,
                ),
                gui.Push(),
            ],

            [gui.Push(),self.collapse(powerpointfile, 'sec_1', True),gui.Push()],
            [gui.Push(),self.collapse(pngfile, 'sec_2', False),gui.Push()],
            [gui.Push(),self.collapse(jpgfile, 'sec_3', False),gui.Push()],
            [gui.Push(),self.collapse(svgfile, 'sec_4', False),gui.Push()],
            
            [
                gui.Push(),
                gui.Button("Cancel"),
                gui.Button("Ok"),
                gui.Push(),
            ],
        ]

        layout = [
            [
                gui.Column(pdf_column),
            ]
        ]
        return layout

    def sudokugen(self):

        
        bpicon_base64 = b'iVBORw0KGgoAAAANSUhEUgAAALQAAAC0CAQAAACXxM65AAAAAXNSR0IArs4c6QAADTlJREFUeNrtnXtwVNUdxz93d/NCIogib0HFFyCgbVF81AmgouPY+KigWKn1MbZq67uiVqctbdXBTu1ofRU66JShttpMdaytiaAI+GpF0GgFBF8xoCJJINkku3v6B2fPbpbdvffc3b27uXt++8ee5N7sufvJ9/7O77x+1yKdWdQykgnMYiqHMogqjGWzbtrYzDoaaaaFDkQ6pHtbNfOpZzLDCRiGWhajlfU0sJSw3alB6mhFmFdOr1bqCKaCTbaR/Jx72N9IM0cbyDkMZQMd6Q+PZQ1Ro8c8vaKsYWw6RY/hCU5K67ONuTGLMRxLE+2kOI1XjAoL8HqFkcmKDvJzzjdqLoCNJkQjIh7e1dHAvoZKQaydelbsAV3NVoYZIgWzbYwjHMBivsFcUBvGfKwAtdQbFgW2emoDjGSyIVFgm8zIAJMYbkgU2IYzKcBMM3RUcAswM8AUw8EDm2KxjQMNh4LbdouwGdb3wLqtdLMBxgrhpo0Z0Aa0MQPagDagjRnQBrQxA9qANqCNeWAhP32ZKk5nIrvZyYe8ab/O0IB2e3Mu4DYq5E/dbOB5VtJMa0mMm3k8emdxDOMJ08F2NtGd188+i78wIOV3nbzPn1lMWwmg9nCBlCVuE90iblHRLO4TJ4phIpSXT79NZLJ1YrqoLPbiMC8rO0507AWhV7wr/iBmiGDeQLeJnaI3pZYdYpHYv3xAr8iouYhYJobk+OkL5H3yExEU1WK6WCi29KnjmXIBPVpERTbbKObkpOubJOhr1W+qxCniIbFN1dAkDioaaA/j6Kk2Qft4HuNuFTXoW49qcFGRx0tcx0n8W7b4ddxPrf87LIcmlbdyFUPZh4OZyzK2EgWglmu43vUl7ZaY+/59NxuZy0P0ABbf4aaiRbSe3TyL1C3cIY4Xlvp9SEwQ94kedewGl59/gfyEG9McqxELRJcQQoiwmOdvH22JRxXoJWmO14sd8miPmJH0b3D+mi3//ta0R2vEUnn8v2Kkn310iBpZivFCmuMNzOFjACq4i0EuauiU7xVpty50cSdbAZjIuX720UEqZaldAk21RuZLX30yl7iooUuNeKS3j/g+MaCSuxjoX9ABFU+EMwz3CFZyI12AxUIO164hrBSdyVbzHAAHcIfnG3YC3lcVk7pNZ0t4XsYfV2tf2m75XpkRYoQGGQTOTd4C6DfQMaXczNbBbyWKU7VR7LRVNKzgSwBGMMuvoIXDYcI1LAHgECZq1tAua6jMcs4W2RBXcrrHzqMIirY762d0AVWcofn5EdkcVmdBKFgsSzNy6IP6QtGwkxcBOFNbc3HQ2WwtWwAYwtxyBx3lNQDGcrAr0FVZ/0ERGXnAD/zuOiwbpQo+IAxY1LkK8GpsPv8N+S8/wdPNwkVQtGXrEr6gF4BvuQI9wOasD2UajQpP9/2VoKKhU0baR7hyHXaK3i5DPLTjGt8pOibPHaW5vSYOOru1qanag/wJOuoYdIW8rEr2KYCi29klS6M8HJv2EHTcdQRT8zrtZfG4IVAQRXfxtSztbxMK9nNF24MeLDsTway9vHS+3a7DsudK4j66xvZK+rmi7W7YAyVgobm6p9OR60gefvIl6IhDRVuMkGdE1RizjuuotHUInSltgc/CO6egA6pHqAs6DnCIo3gbAh4OLHkIutcx6Ph8eYdCp6No2M/mvLjn78kyMu4L0Nl9dJDxsvQ/hyN+qUq1S0VZpa4o5mfQlk0sMYmhsvSuqy64veuoUef7WtHYRMfzlOd8vUCgh6muS48fQSe+VChrr/C7KghbqVlDtyPQNRwgS1/6E3TEEehjGCVLazVjDqeK3lctNmjxcBV+UVxHKEsMfaYqN2rXEFY9y2w2REUlW7378sUBXZFFbSfI0g5e0q4h7joG2fQ7D3DZ2PrIdYzgG7K0kU0uaog5UPTBculuhPf9CdpJY3iR8q+r+Eq7hvjSnOyKPkW+v5LnzUolCDqT66jlh6rzvdhFQxV1oOigWsawHPwJuttG0RaXKe/5V1e3dVzR2brgs2QU3cPf/Ao6bAP6QL4nS7tY5KqGmFT0vhkHi4JcI0sb1PB/2YE+iQmy9ALvuaohKhUdzDjLchTTZWmNx/tpSwZ0LbfKceQISzRH7VIVnWk6y+I02VB2stK/oLP76F+qwG6VXLqbC+jqDN37ebLur1iLt+bhFqXM4Z3F2Vwp/ep2LkuKuPUsMfWVHvTFHCtL/+Rzj0EXpcOSGt6N41fyZo/we7kI0R3obIoex73Kid0L/lV0NG2lFuNYqZayfMDjOdSQUPTeI96V3KImBJax2XPQRe+CT2W5wtzN7XySQw2xLKCnq+HXr3P6Z/YD0NG9XEeQK1jDNPX7+2jIqYbMih7IE6oz9CZvFgF0qBiKHs23GcThnMe0pInatTyYYw0JH913Dmc/HmWMapLvVus6fA/6ZBqoZEBK/+2BnCOBTIq+krNVebH2vE2/jjr2Y5+9usk/ynlLWgJ0RVJdl7FQgd/EQg9nvovuo9PbdBar+e983agBLuJ36qcoi2gFv4Putb2UGTxlO4OtA9qijgeStiM38sci6bmkFL1H1b/J28LDALN5Kgnzf7jYw3UcJdEYJqyLJt5gKJfLvlyIy+nl2jwM+AT5MXcmzbXs4Ga1XLc45mHOipjMmBEWm0STuFvMEDUyl8cU8bhoU4msrnKZW2mS+FjlZ/qsT76mDnG1CJRPlrAu+bWfE4eIwSnHqsS5okUe/1RMzBF0X9slrshTbr1+kbwqMVAa5kO1RT5x7Gnq2YQARvFkHh+x+AW38ifXI4L9sDFMeOlMa+9e53zeAWACC233Czrzi6s4jweLjrlIoDOvJn2bm+VZF2nvm03dihGlnV9zGqtKIglsSSka4F/8lAgwgEXaT0tLgP6cZdzAVO4ooeTGHjYIH8nG6bWsZ1WJZ+V5izWzhR2lkmQuFVUl9oTwojSGlTZnPSTH187hZJeKjtn2Q33tOsKOQEOTnDodzKVa23kSoIMl9zT5ooCusj1vj5+2mMMhrkAHDGgnoOEtXgaghkc0LtEoWnlOZ64DBEvk2TM5WgN0TH2tQDmDjpuTtFHr+UyW6h1fZMwouq852Uq/RW18ONFxzudY0sb+sgYtNBS9i1WydLSav7YHHTONISTmWCocYVgh34c5zq2UrOiAAQ2WI02/Ks+3mK3to8tc0REtL72LV2VppvZlGh/tOJIGWC3fR/d5PoCzr2WiDseKhjdU6UijaHeVOVP0u2pN9RHaNZV1Y1ihCXo3X8jSYUbR7kA7Syncww5ZGm58tDvPHHIIOr5FbbDmApRQebsOXdC9aq68yjZtYOq9Y3y0BuiISm4Z0sxfXuFhRjsfKDqRji2oCa7CuA53oAOaoCvLW9EVSeGXs55kpxodsYyiC6doNJNLJYZJK4yi9UBHNEFHjaLdKbpXORGhpehKE3XogY5KwFFH2jaKdu064km3nYGOJvnoUPmCtpK+vFMMA5Wv7tVyHQEPk8qXHOjkEbWgw4sbqKKPbi3XgeZjGHwGGk1FV6uFu50OfXTMgO4L1ynoeJ7SnZpxtFG0FugBjJYlZ7vEo0muo8aAdu6jj2CELG3UVvSA8gVdpa3oC9Tlve8QdNSA7pvnyAnoQer53RHe0h4bqTGKdgp6jlpzt8FhY5i8sd8o2qGP3pc5qvys4zq609ZWZqArtar9plra2MvfXYCuMqCdjYrcolZFN9NiFI0WPOd2edID0tdqpOyOGNA6ip7GQuXFYzRoDP9HjetwDrqGXySl2n5BLUjXU3RliYEOlRroaho4Tf3UzqVa84aiZEEXKY7O3Em5PwlzD/e4zupVYRSdLXZentQIwmoecJ0Cwig6gwWZxQpmJ+Woa+I82vuDgkrterJVNZ7bOavPNrdnuDanxx1YBnQyjBpqGcMcrujzkJoOnuSaEkpp0s9AJ3upMxEEGcJoxnFkykjbZm7n6ZLLt9GPQCffzMdxXIYOx5NckkOqqYgBbW9hVvMwz+YEK9F47jag09s7LGAVbTl+SiIpZkv5gg5n6Mvt4D0eYXlebvuN6lPXlBhoy7ukcAexlFP6eOpPeJFVNNOcs5ITHfj1HAZs5IwiPJmiREADHM+pHEYtITbzD14uQOM1iQVs41FPHwpZgqDL1wIGgQFtQBszoA1oA9qYAW1AGzOgDWgD2pg3oLsNBA+sO5C3EUpj2awtUGLDtn61zQHWGQoe2LoAjUV7lmL5WIzGAM1Fezpo+VgrzQFaWG9IFNjW0xKgI8ennhuztwY6LKCarXl8fKCxVNvGOMIBIMyFOayPNZbd2rmQcHxn5ScMZVrJrXT1R7zxMI8h4qAFGziRMYZL3u1VrtvjLeK7zDpo5FjGGjJ5tZe5mE/3FBObsttoYhqjjQPJm9NYwzw+jv8Y7OO2nyfEhJLbC9k/m8CHuV497iuNBamjFWFeOb1aqUtN4JDOUVQzn3omM9xMC2i7i1bW08DSvdcop/fIFrWMZBIzmcKhDDLOxMa6aWMzb9PEO7TQkW7h6P8BBU+G2vkWH04AAAAASUVORK5CYII#'


        LICNO = """ASPOSE LICENSE"""

        global presentation
        p = os.path.expanduser('~/SudokuGen/')
        try:
            if os.path.exists(p) == False:
                os.makedirs(p)
        except OSError as err:
                                print(err)
        pass

        filepath = os.path.expanduser('~/SudokuGen/license.lic')

        if os.path.exists(filepath) == False:
            text_file = open(filepath, "w")
            n = text_file.write(LICNO)
            text_file.close()

        license = slides2.License()
        license.set_license(filepath);

        output, success = self.system_call(exportfile+' --generate ' + str(noofpuzzles) + ' --difficulty ' + str(difficulty) + ' --solution ' + '--csv ' + '--one-line')

        if success == True:
            with open(p + 'output.txt', 'w') as f:
                f.write(output)
            sudokulocation = p + 'output.txt'
            presentation = Presentation()
            presentation.slide_width = Inches(12)
            presentation.slide_height = Inches(12)

            fileObj = [line for line in open(sudokulocation, "r").read().splitlines() if line]
            words = fileObj

            array_length = len(words)

            s = 0

            layout2 = [[gui.Text('Creating your Sudoku Puzzles...', key='updatetext')],
                        [gui.ProgressBar(max_value=int(array_length), orientation='h', size=(40, 20), key='progress')]]

            window2 = gui.Window('Creating your Sudoku Puzzles...', layout2, finalize=True, background_color='#000000', icon=(bpicon_base64))
                    
            for z in range(array_length):
                if PPTX == True:
                    s = s + 1
                    progress_bar = window2['progress']
                    progress_bar.update_bar(s)

                currentpuzzle = words[z]
                layout = presentation.slide_masters[0].slide_layouts[6]
                slide = presentation.slides.add_slide(layout)
                Top = (0 + Inches(1))
                Left = (0 + Inches(1))
                Width = (presentation.slide_width - (Left * 2))
                Height = (presentation.slide_height - (Top * 2))
                SudokuW = 9
                SudokuH = 9
                           

                if (z % 2) == 0:
                    title = str(sudokutitle) + " "
                    number = str(int((z/2) + 1))
                else:
                    title = str(sudokusolutiontitle) + " "
                    number = str(int((z/2) + 0.5))
                if z==0:
                    title = str(sudokutitle) + " "
                    number = str(z+1)
                if z==1:
                    title = str(sudokusolutiontitle) + " "
                    number = str(z)

                txBox = slide.shapes.add_textbox((0 + Inches(1)), (0 + Inches(0.1)), (presentation.slide_width - (Left * 2)), Inches(0.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = title+number
                #p.font.name = TEMPLATEFONT
                p.font.size = Pt(44)
                p.alignment = PP_ALIGN.CENTER


                
                shape = slide.shapes.add_table(SudokuW, SudokuH, Top, Left, Width, Height)

                tbl =  shape._element.graphic.graphicData.tbl
                style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'

                tbl[0][-1].text = style_id

                c = 0
                for i in range(SudokuW):
                    for j in range(SudokuH):
                        if currentpuzzle[c] == ".":
                            shape.table.cell(i,j).text == " "
                            if j==0:
                                cell = shape.table.cell(i,j)
                                cell = self._set_left_cell_border(cell)
                            elif j==2:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            elif j==5:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            elif j==8:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            if i ==0:
                                cell = shape.table.cell(i,j)
                                cell = self._set_top_cell_border(cell)
                            elif i==2:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            elif i==5:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            elif i==8:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            
                            
                                
                        else:
                            shape.table.cell(i,j).text = currentpuzzle[c]
                            shape.table.cell(i,j).text_frame.paragraphs[0].font.size = Pt(54)
                            shape.table.cell(i,j).text_frame.paragraphs[0].margin_bottom = 0
                            shape.table.cell(i,j).text_frame.paragraphs[0].margin_top = 0
                            shape.table.cell(i,j).vertical_anchor = MSO_ANCHOR.MIDDLE
                            shape.table.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                            if j==0:
                                cell = shape.table.cell(i,j)
                                cell = self._set_left_cell_border(cell)
                            elif j==2:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            elif j==5:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            elif j==8:
                                cell = shape.table.cell(i,j)
                                cell = self._set_right_cell_border(cell)
                            if i==0:
                                cell = shape.table.cell(i,j)
                                cell = self._set_top_cell_border(cell)
                            elif i==2:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            elif i==5:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            elif i==8:
                                cell = shape.table.cell(i,j)
                                cell = self._set_bottom_cell_border(cell)
                            
                            

                            
                        c = c + 1
            for b in reversed(range(array_length)):
                if (b % 2) != 0:
                    self.move_slide(b, array_length)

            for t in reversed(range(array_length)):
                if t >= (array_length / 2):
                    self.move_slide(t, array_length)

            filepath2 = os.path.expanduser("~/SudokuGen/temp.pptx")

            presentation.save(filepath2)
            presentation2 = slides2.Presentation(filepath2)

            my_os=sys.platform
            loc = ""
            if my_os == "darwin":
                loc = "//"
            else:
                loc = "\\" 
            if my_os == "linux":
                loc = "//"

            if PPTX == True:
                presentation2.save(exportfile2, slides2.export.SaveFormat.PPTX)

            if PNG == True:
                for index in range(presentation2.slides.length):
                    s = s + 1
                    progress_bar = window2['progress']
                    progress_bar.update_bar(s)
                    no = int(index) + 1
                    slide = presentation2.slides[index]
                    if no <= array_length / 2:
                        slide.get_thumbnail(300/72,300/72).save(exportfile2 + loc + "sudoku_{no}.png".format(no=no), drawing.imaging.ImageFormat.png)
                    else:
                        no = int(no - (array_length / 2))
                        slide.get_thumbnail(300/72,300/72).save(exportfile3 + loc + "sudoku_solution_{no}.png".format(no=no), drawing.imaging.ImageFormat.png)        
            if JPG == True:
                for index in range(presentation2.slides.length):
                    s = s + 1
                    progress_bar = window2['progress']
                    progress_bar.update_bar(s)
                    no = int(index) + 1
                    slide = presentation2.slides[index]
                    if no <= array_length / 2:
                        slide.get_thumbnail(300/72,300/72).save(exportfile2 + loc + "sudoku_{no}.jpg".format(no=no), drawing.imaging.ImageFormat.jpeg)
                    else:
                        no = int(no - (array_length / 2))
                        slide.get_thumbnail(300/72,300/72).save(exportfile3 + loc + "sudoku_solution_{no}.jpg".format(no=no), drawing.imaging.ImageFormat.jpeg)          
            if SVG == True:
                for index in range(presentation2.slides.length):
                    s = s + 1
                    progress_bar = window2['progress']
                    progress_bar.update_bar(s)
                    no = int(index) + 1
                    slide = presentation2.slides[index]
                    if no <= array_length / 2:
                        with open(exportfile2 + loc + "sudoku_{no}.svg".format(no=no), "wb") as file:
                            slide.write_as_svg(file)
                    else:
                        no = int(no - (array_length / 2))
                        with open(exportfile3 + loc + "sudoku_solution_{no}.svg".format(no=no), "wb") as file:
                            slide.write_as_svg(file)

            window2.close()
        else:
            gui.Popup(output)

